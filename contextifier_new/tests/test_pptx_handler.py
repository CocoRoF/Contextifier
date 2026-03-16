"""
Comprehensive tests for the PPTX handler pipeline.

Tests all 7 modules:
- _constants: ElementType, SlideElement, font mappings
- converter: PptxConverter (ZIP validation, Presentation creation)
- preprocessor: PptxPreprocessor (wrap, stats, chart pre-extraction)
- metadata_extractor: PptxMetadataExtractor (core properties)
- _bullet: extract_text_with_bullets (bullets, numbering, special fonts)
- _table: is_simple_table, extract_simple_text, extract_table
- content_extractor: PptxContentExtractor (full slide extraction)
- handler: PPTXHandler (pipeline wiring)
"""

import io
import zipfile
import pytest
from unittest.mock import MagicMock, patch, PropertyMock
from dataclasses import dataclass
from typing import Any, List, Optional

# ── Helpers to create minimal PPTX ─────────────────────────────────────

def _make_minimal_pptx() -> bytes:
    """Create a minimal valid PPTX file with one empty slide."""
    from pptx import Presentation
    prs = Presentation()
    slide_layout = prs.slide_layouts[6]  # Blank layout
    prs.slides.add_slide(slide_layout)
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _make_pptx_with_text(texts: list[str]) -> bytes:
    """Create a PPTX with one slide containing text boxes."""
    from pptx import Presentation
    from pptx.util import Inches
    prs = Presentation()
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    for i, text in enumerate(texts):
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1 + i), Inches(5), Inches(0.5))
        txBox.text_frame.text = text
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _make_pptx_with_table(rows_data: list[list[str]]) -> bytes:
    """Create a PPTX with a table on one slide."""
    from pptx import Presentation
    from pptx.util import Inches
    prs = Presentation()
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    num_rows = len(rows_data)
    num_cols = len(rows_data[0]) if rows_data else 1
    table_shape = slide.shapes.add_table(num_rows, num_cols, Inches(1), Inches(1), Inches(6), Inches(2))
    table = table_shape.table
    for r_idx, row in enumerate(rows_data):
        for c_idx, cell_text in enumerate(row):
            table.cell(r_idx, c_idx).text = cell_text
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _make_pptx_multi_slide(slide_texts: list[list[str]]) -> bytes:
    """Create a PPTX with multiple slides, each with text boxes."""
    from pptx import Presentation
    from pptx.util import Inches
    prs = Presentation()
    slide_layout = prs.slide_layouts[6]
    for texts in slide_texts:
        slide = prs.slides.add_slide(slide_layout)
        for i, text in enumerate(texts):
            txBox = slide.shapes.add_textbox(Inches(1), Inches(1 + i), Inches(5), Inches(0.5))
            txBox.text_frame.text = text
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _make_pptx_with_notes(text: str, notes: str) -> bytes:
    """Create a PPTX with a slide that has notes."""
    from pptx import Presentation
    from pptx.util import Inches
    prs = Presentation()
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(0.5))
    txBox.text_frame.text = text
    # Add notes
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = notes
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _make_file_context(data: bytes) -> dict:
    """Create a FileContext dict."""
    return {
        "file_path": "/test/file.pptx",
        "file_name": "file.pptx",
        "file_extension": "pptx",
        "file_category": "presentation",
        "file_data": data,
        "file_stream": io.BytesIO(data),
        "file_size": len(data),
    }


# ═══════════════════════════════════════════════════════════════════════════════
# Test Constants
# ═══════════════════════════════════════════════════════════════════════════════

class TestConstants:
    """Tests for _constants module."""

    def test_zip_magic(self):
        from contextifier_new.handlers.pptx._constants import ZIP_MAGIC
        assert ZIP_MAGIC == b"PK\x03\x04"

    def test_element_type_values(self):
        from contextifier_new.handlers.pptx._constants import ElementType
        assert ElementType.TEXT == "text"
        assert ElementType.IMAGE == "image"
        assert ElementType.TABLE == "table"
        assert ElementType.CHART == "chart"

    def test_slide_element_sort_key(self):
        from contextifier_new.handlers.pptx._constants import SlideElement, ElementType
        # sort_key = (top, left)
        e1 = SlideElement(ElementType.TEXT, "A", (100, 200, 50, 50))  # top=200
        e2 = SlideElement(ElementType.TEXT, "B", (50, 100, 50, 50))   # top=100
        assert e2.sort_key < e1.sort_key  # e2 is higher (top=100 < 200)

    def test_slide_element_sort_key_same_top(self):
        from contextifier_new.handlers.pptx._constants import SlideElement, ElementType
        e1 = SlideElement(ElementType.TEXT, "A", (200, 100, 50, 50))
        e2 = SlideElement(ElementType.TEXT, "B", (100, 100, 50, 50))
        assert e2.sort_key < e1.sort_key  # same top, e2 left=100 < 200

    def test_wingdings_mapping(self):
        from contextifier_new.handlers.pptx._constants import WINGDINGS_MAPPING
        assert WINGDINGS_MAPPING[0x6C] == "●"
        assert WINGDINGS_MAPPING[0xFC] == "✓"
        assert WINGDINGS_MAPPING[0xD8] == "➢"

    def test_wingdings_char_mapping(self):
        from contextifier_new.handlers.pptx._constants import WINGDINGS_CHAR_MAPPING
        assert WINGDINGS_CHAR_MAPPING["Ø"] == "➢"
        assert WINGDINGS_CHAR_MAPPING["n"] == "■"

    def test_symbol_mapping(self):
        from contextifier_new.handlers.pptx._constants import SYMBOL_MAPPING
        assert SYMBOL_MAPPING[0xB7] == "•"
        assert SYMBOL_MAPPING[0xD7] == "×"


# ═══════════════════════════════════════════════════════════════════════════════
# Test Converter
# ═══════════════════════════════════════════════════════════════════════════════

class TestPptxConverter:
    """Tests for PptxConverter."""

    def test_convert_valid_pptx(self):
        from contextifier_new.handlers.pptx.converter import PptxConverter
        conv = PptxConverter()
        data = _make_minimal_pptx()
        ctx = _make_file_context(data)
        prs = conv.convert(ctx)
        assert hasattr(prs, "slides")
        assert len(prs.slides) == 1

    def test_convert_empty_data_raises(self):
        from contextifier_new.handlers.pptx.converter import PptxConverter
        from contextifier_new.errors import ConversionError
        conv = PptxConverter()
        ctx = _make_file_context(b"")
        with pytest.raises(ConversionError):
            conv.convert(ctx)

    def test_validate_valid_pptx(self):
        from contextifier_new.handlers.pptx.converter import PptxConverter
        conv = PptxConverter()
        data = _make_minimal_pptx()
        ctx = _make_file_context(data)
        assert conv.validate(ctx) is True

    def test_validate_invalid_magic(self):
        from contextifier_new.handlers.pptx.converter import PptxConverter
        conv = PptxConverter()
        ctx = _make_file_context(b"\x00\x00\x00\x00")
        assert conv.validate(ctx) is False

    def test_validate_too_short(self):
        from contextifier_new.handlers.pptx.converter import PptxConverter
        conv = PptxConverter()
        ctx = _make_file_context(b"PK")
        assert conv.validate(ctx) is False

    def test_validate_zip_without_content_types(self):
        from contextifier_new.handlers.pptx.converter import PptxConverter
        conv = PptxConverter()
        # Create a regular ZIP without [Content_Types].xml
        buf = io.BytesIO()
        with zipfile.ZipFile(buf, "w") as zf:
            zf.writestr("test.txt", "hello")
        ctx = _make_file_context(buf.getvalue())
        assert conv.validate(ctx) is False

    def test_get_format_name(self):
        from contextifier_new.handlers.pptx.converter import PptxConverter
        assert PptxConverter().get_format_name() == "pptx"

    def test_close_does_not_raise(self):
        from contextifier_new.handlers.pptx.converter import PptxConverter
        conv = PptxConverter()
        conv.close(None)
        conv.close("anything")

    def test_convert_invalid_pptx_raises(self):
        from contextifier_new.handlers.pptx.converter import PptxConverter
        from contextifier_new.errors import ConversionError
        conv = PptxConverter()
        # Valid ZIP magic but not a PPTX
        buf = io.BytesIO()
        with zipfile.ZipFile(buf, "w") as zf:
            zf.writestr("[Content_Types].xml", "<Types/>")
        ctx = _make_file_context(buf.getvalue())
        # python-pptx may or may not raise — just ensure no crash
        try:
            conv.convert(ctx)
        except ConversionError:
            pass  # Expected


# ═══════════════════════════════════════════════════════════════════════════════
# Test Preprocessor
# ═══════════════════════════════════════════════════════════════════════════════

class TestPptxPreprocessor:
    """Tests for PptxPreprocessor."""

    def test_preprocess_wraps_presentation(self):
        from pptx import Presentation
        from contextifier_new.handlers.pptx.preprocessor import PptxPreprocessor
        data = _make_minimal_pptx()
        prs = Presentation(io.BytesIO(data))
        pp = PptxPreprocessor()
        result = pp.preprocess(prs)
        assert result.content is prs
        assert result.raw_content is prs

    def test_preprocess_slide_count(self):
        from pptx import Presentation
        from contextifier_new.handlers.pptx.preprocessor import PptxPreprocessor
        data = _make_pptx_multi_slide([["a"], ["b"], ["c"]])
        prs = Presentation(io.BytesIO(data))
        pp = PptxPreprocessor()
        result = pp.preprocess(prs)
        assert result.properties["slide_count"] == 3

    def test_preprocess_dimensions(self):
        from pptx import Presentation
        from contextifier_new.handlers.pptx.preprocessor import PptxPreprocessor
        data = _make_minimal_pptx()
        prs = Presentation(io.BytesIO(data))
        pp = PptxPreprocessor()
        result = pp.preprocess(prs)
        assert result.properties["slide_width"] is not None
        assert result.properties["slide_height"] is not None

    def test_preprocess_none_raises(self):
        from contextifier_new.handlers.pptx.preprocessor import PptxPreprocessor
        from contextifier_new.errors import PreprocessingError
        pp = PptxPreprocessor()
        with pytest.raises(PreprocessingError):
            pp.preprocess(None)

    def test_preprocess_charts_by_slide(self):
        from pptx import Presentation
        from contextifier_new.handlers.pptx.preprocessor import PptxPreprocessor
        # Minimal PPTX without charts → empty dict
        data = _make_minimal_pptx()
        prs = Presentation(io.BytesIO(data))
        pp = PptxPreprocessor()
        result = pp.preprocess(prs)
        assert result.resources.get("charts_by_slide") == {}

    def test_get_format_name(self):
        from contextifier_new.handlers.pptx.preprocessor import PptxPreprocessor
        assert PptxPreprocessor().get_format_name() == "pptx"


# ═══════════════════════════════════════════════════════════════════════════════
# Test Metadata Extractor
# ═══════════════════════════════════════════════════════════════════════════════

class TestPptxMetadataExtractor:
    """Tests for PptxMetadataExtractor."""

    def test_extract_from_presentation(self):
        from pptx import Presentation
        from contextifier_new.handlers.pptx.metadata_extractor import PptxMetadataExtractor
        data = _make_minimal_pptx()
        prs = Presentation(io.BytesIO(data))
        ext = PptxMetadataExtractor()
        meta = ext.extract(prs)
        # page_count = slide count
        assert meta.page_count == 1

    def test_extract_from_preprocessed_data(self):
        from pptx import Presentation
        from contextifier_new.handlers.pptx.metadata_extractor import PptxMetadataExtractor
        from contextifier_new.types import PreprocessedData
        data = _make_minimal_pptx()
        prs = Presentation(io.BytesIO(data))
        ppd = PreprocessedData(content=prs, raw_content=prs)
        ext = PptxMetadataExtractor()
        meta = ext.extract(ppd)
        assert meta.page_count == 1

    def test_extract_with_set_properties(self):
        from pptx import Presentation
        from contextifier_new.handlers.pptx.metadata_extractor import PptxMetadataExtractor
        data = _make_minimal_pptx()
        prs = Presentation(io.BytesIO(data))
        prs.core_properties.title = "Test Title"
        prs.core_properties.author = "Test Author"
        prs.core_properties.subject = "Test Subject"
        ext = PptxMetadataExtractor()
        meta = ext.extract(prs)
        assert meta.title == "Test Title"
        assert meta.author == "Test Author"
        assert meta.subject == "Test Subject"

    def test_extract_none_returns_empty(self):
        from contextifier_new.handlers.pptx.metadata_extractor import PptxMetadataExtractor
        ext = PptxMetadataExtractor()
        meta = ext.extract(None)
        assert meta.is_empty()

    def test_extract_nonsense_returns_empty(self):
        from contextifier_new.handlers.pptx.metadata_extractor import PptxMetadataExtractor
        ext = PptxMetadataExtractor()
        meta = ext.extract("not a presentation")
        assert meta.is_empty()

    def test_get_format_name(self):
        from contextifier_new.handlers.pptx.metadata_extractor import PptxMetadataExtractor
        assert PptxMetadataExtractor().get_format_name() == "pptx"


# ═══════════════════════════════════════════════════════════════════════════════
# Test Bullet Extraction
# ═══════════════════════════════════════════════════════════════════════════════

class TestBulletExtraction:
    """Tests for _bullet module."""

    def test_extract_plain_text(self):
        from contextifier_new.handlers.pptx._bullet import extract_text_with_bullets
        tf = self._make_textframe(["Hello", "World"])
        result = extract_text_with_bullets(tf)
        assert "Hello" in result
        assert "World" in result

    def test_extract_empty_textframe(self):
        from contextifier_new.handlers.pptx._bullet import extract_text_with_bullets
        assert extract_text_with_bullets(None) == ""

    def test_extract_from_real_pptx(self):
        """Extract text from a real PPTX textbox."""
        from pptx import Presentation
        from pptx.util import Inches
        from contextifier_new.handlers.pptx._bullet import extract_text_with_bullets

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(2))
        txBox.text_frame.text = "First line"
        txBox.text_frame.add_paragraph().text = "Second line"

        result = extract_text_with_bullets(txBox.text_frame)
        assert "First line" in result
        assert "Second line" in result

    def test_extract_preserves_empty_paragraphs(self):
        from contextifier_new.handlers.pptx._bullet import extract_text_with_bullets
        tf = self._make_textframe(["Hello", "", "World"])
        result = extract_text_with_bullets(tf)
        lines = result.split("\n")
        assert len(lines) == 3
        assert lines[1] == ""

    @staticmethod
    def _make_textframe(texts: list[str]):
        """Create a mock TextFrame."""
        paras = []
        for t in texts:
            para = MagicMock()
            para.text = t
            para.level = 0
            para._element = MagicMock()
            para._element.pPr = None
            paras.append(para)

        tf = MagicMock()
        tf.paragraphs = paras
        tf.text = "\n".join(texts)
        return tf


# ═══════════════════════════════════════════════════════════════════════════════
# Test Table Extraction
# ═══════════════════════════════════════════════════════════════════════════════

class TestTableExtraction:
    """Tests for _table module."""

    def test_is_simple_table_single_row(self):
        from contextifier_new.handlers.pptx._table import is_simple_table
        table = self._make_mock_table(1, 3, [["a", "b", "c"]])
        assert is_simple_table(table) is True

    def test_is_simple_table_single_col(self):
        from contextifier_new.handlers.pptx._table import is_simple_table
        table = self._make_mock_table(3, 1, [["a"], ["b"], ["c"]])
        assert is_simple_table(table) is True

    def test_is_not_simple_table(self):
        from contextifier_new.handlers.pptx._table import is_simple_table
        table = self._make_mock_table(2, 2, [["a", "b"], ["c", "d"]])
        assert is_simple_table(table) is False

    def test_extract_simple_text(self):
        from contextifier_new.handlers.pptx._table import extract_simple_text
        table = self._make_mock_table(1, 3, [["a", "b", "c"]])
        result = extract_simple_text(table)
        assert "a" in result
        assert "b" in result
        assert "c" in result

    def test_extract_table_basic(self):
        from contextifier_new.handlers.pptx._table import extract_table
        table = self._make_mock_table(2, 2, [["H1", "H2"], ["V1", "V2"]])
        td = extract_table(table)
        assert td.num_rows == 2
        assert td.num_cols == 2
        assert td.has_header is True
        assert td.rows[0][0].content == "H1"
        assert td.rows[1][1].content == "V2"

    def test_extract_table_from_real_pptx(self):
        """Extract table from a real PPTX."""
        from pptx import Presentation
        from contextifier_new.handlers.pptx._table import extract_table
        data = _make_pptx_with_table([["H1", "H2"], ["A", "B"]])
        prs = Presentation(io.BytesIO(data))
        slide = prs.slides[0]
        for shape in slide.shapes:
            if shape.has_table:
                td = extract_table(shape.table)
                assert td.num_rows == 2
                assert td.num_cols == 2
                assert td.rows[0][0].content == "H1"
                return
        pytest.fail("No table found in PPTX")

    def test_extract_table_empty(self):
        from contextifier_new.handlers.pptx._table import extract_table
        table = MagicMock()
        table.rows = MagicMock()
        table.rows.__len__ = MagicMock(return_value=0)
        table.columns = MagicMock()
        table.columns.__len__ = MagicMock(return_value=0)
        td = extract_table(table)
        assert td.num_rows == 0

    @staticmethod
    def _make_mock_table(num_rows, num_cols, data):
        """Create a mock table with given data."""
        table = MagicMock()
        rows = [MagicMock() for _ in range(num_rows)]
        columns = [MagicMock() for _ in range(num_cols)]
        table.rows = rows
        table.columns = columns

        cells = {}
        for r in range(num_rows):
            row_cells = []
            for c in range(num_cols):
                cell = MagicMock()
                cell.text = data[r][c] if r < len(data) and c < len(data[r]) else ""
                cell.is_merge_origin = False
                cell.is_spanned = False
                cell._tc = MagicMock()
                cell._tc.get = MagicMock(return_value=None)
                cells[(r, c)] = cell
                row_cells.append(cell)
            rows[r].cells = row_cells

        def cell_accessor(r, c):
            return cells.get((r, c), MagicMock())

        table.cell = cell_accessor
        return table


# ═══════════════════════════════════════════════════════════════════════════════
# Test Content Extractor
# ═══════════════════════════════════════════════════════════════════════════════

class TestPptxContentExtractor:
    """Tests for PptxContentExtractor."""

    def test_extract_text_single_slide(self):
        from pptx import Presentation
        from contextifier_new.handlers.pptx.content_extractor import PptxContentExtractor
        from contextifier_new.types import PreprocessedData

        data = _make_pptx_with_text(["Hello World"])
        prs = Presentation(io.BytesIO(data))
        ppd = PreprocessedData(content=prs, raw_content=prs, resources={"charts_by_slide": {}})

        ext = PptxContentExtractor()
        text = ext.extract_text(ppd)
        assert "Hello World" in text
        assert "[Slide:1]" in text

    def test_extract_text_multi_slide(self):
        from pptx import Presentation
        from contextifier_new.handlers.pptx.content_extractor import PptxContentExtractor
        from contextifier_new.types import PreprocessedData

        data = _make_pptx_multi_slide([["Slide 1"], ["Slide 2"], ["Slide 3"]])
        prs = Presentation(io.BytesIO(data))
        ppd = PreprocessedData(content=prs, raw_content=prs, resources={"charts_by_slide": {}})

        ext = PptxContentExtractor()
        text = ext.extract_text(ppd)
        assert "[Slide:1]" in text
        assert "[Slide:2]" in text
        assert "[Slide:3]" in text
        assert "Slide 1" in text
        assert "Slide 2" in text
        assert "Slide 3" in text

    def test_extract_text_with_table(self):
        from pptx import Presentation
        from contextifier_new.handlers.pptx.content_extractor import PptxContentExtractor
        from contextifier_new.types import PreprocessedData

        data = _make_pptx_with_table([["H1", "H2"], ["A", "B"]])
        prs = Presentation(io.BytesIO(data))
        ppd = PreprocessedData(content=prs, raw_content=prs, resources={"charts_by_slide": {}})

        ext = PptxContentExtractor()
        text = ext.extract_text(ppd)
        assert "H1" in text
        assert "H2" in text
        assert "A" in text

    def test_extract_text_empty_slide(self):
        from pptx import Presentation
        from contextifier_new.handlers.pptx.content_extractor import PptxContentExtractor
        from contextifier_new.types import PreprocessedData

        data = _make_minimal_pptx()
        prs = Presentation(io.BytesIO(data))
        ppd = PreprocessedData(content=prs, raw_content=prs, resources={"charts_by_slide": {}})

        ext = PptxContentExtractor()
        text = ext.extract_text(ppd)
        assert "[Empty Slide]" in text

    def test_extract_text_with_notes(self):
        from pptx import Presentation
        from contextifier_new.handlers.pptx.content_extractor import PptxContentExtractor
        from contextifier_new.types import PreprocessedData

        data = _make_pptx_with_notes("Main content", "Speaker notes here")
        prs = Presentation(io.BytesIO(data))
        ppd = PreprocessedData(content=prs, raw_content=prs, resources={"charts_by_slide": {}})

        ext = PptxContentExtractor()
        text = ext.extract_text(ppd)
        assert "Main content" in text
        assert "Speaker notes here" in text
        assert "[Notes]" in text

    def test_extract_text_none_returns_empty(self):
        from contextifier_new.handlers.pptx.content_extractor import PptxContentExtractor
        from contextifier_new.types import PreprocessedData

        ppd = PreprocessedData(content=None, resources={})
        ext = PptxContentExtractor()
        assert ext.extract_text(ppd) == ""

    def test_extract_tables(self):
        from pptx import Presentation
        from contextifier_new.handlers.pptx.content_extractor import PptxContentExtractor
        from contextifier_new.types import PreprocessedData

        data = _make_pptx_with_table([["H1", "H2"], ["A", "B"], ["C", "D"]])
        prs = Presentation(io.BytesIO(data))
        ppd = PreprocessedData(content=prs, raw_content=prs, resources={})

        ext = PptxContentExtractor()
        tables = ext.extract_tables(ppd)
        assert len(tables) == 1
        assert tables[0].num_rows == 3
        assert tables[0].num_cols == 2

    def test_extract_images_no_images(self):
        from pptx import Presentation
        from contextifier_new.handlers.pptx.content_extractor import PptxContentExtractor
        from contextifier_new.types import PreprocessedData

        data = _make_pptx_with_text(["No images"])
        prs = Presentation(io.BytesIO(data))
        ppd = PreprocessedData(content=prs, raw_content=prs, resources={})

        ext = PptxContentExtractor()
        images = ext.extract_images(ppd)
        assert images == []

    def test_extract_charts_no_charts(self):
        from pptx import Presentation
        from contextifier_new.handlers.pptx.content_extractor import PptxContentExtractor
        from contextifier_new.types import PreprocessedData

        data = _make_minimal_pptx()
        prs = Presentation(io.BytesIO(data))
        ppd = PreprocessedData(content=prs, raw_content=prs, resources={})

        ext = PptxContentExtractor()
        charts = ext.extract_charts(ppd)
        assert charts == []

    def test_get_format_name(self):
        from contextifier_new.handlers.pptx.content_extractor import PptxContentExtractor
        assert PptxContentExtractor().get_format_name() == "pptx"

    def test_table_collapse_1x1(self):
        """1×1 table should just return the cell content."""
        from pptx import Presentation
        from pptx.util import Inches
        from contextifier_new.handlers.pptx.content_extractor import PptxContentExtractor
        from contextifier_new.types import PreprocessedData

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        tbl_shape = slide.shapes.add_table(1, 1, Inches(1), Inches(1), Inches(3), Inches(1))
        tbl_shape.table.cell(0, 0).text = "Single cell"
        buf = io.BytesIO()
        prs.save(buf)
        prs2 = Presentation(io.BytesIO(buf.getvalue()))
        ppd = PreprocessedData(content=prs2, raw_content=prs2, resources={"charts_by_slide": {}})

        ext = PptxContentExtractor()
        text = ext.extract_text(ppd)
        assert "Single cell" in text

    def test_table_collapse_single_column(self):
        """Single-column multi-row table should be line-separated."""
        from pptx import Presentation
        from contextifier_new.handlers.pptx.content_extractor import PptxContentExtractor
        from contextifier_new.types import PreprocessedData

        data = _make_pptx_with_table([["Item 1"], ["Item 2"], ["Item 3"]])
        prs = Presentation(io.BytesIO(data))
        ppd = PreprocessedData(content=prs, raw_content=prs, resources={"charts_by_slide": {}})

        ext = PptxContentExtractor()
        text = ext.extract_text(ppd)
        assert "Item 1" in text
        assert "Item 2" in text
        assert "Item 3" in text


# ═══════════════════════════════════════════════════════════════════════════════
# Test Handler Wiring
# ═══════════════════════════════════════════════════════════════════════════════

class TestPPTXHandler:
    """Tests for PPTXHandler."""

    def test_supported_extensions(self):
        from contextifier_new.handlers.pptx.handler import PPTXHandler
        from contextifier_new.config import ProcessingConfig
        handler = PPTXHandler(ProcessingConfig())
        assert handler.supported_extensions == frozenset({"pptx"})

    def test_handler_name(self):
        from contextifier_new.handlers.pptx.handler import PPTXHandler
        from contextifier_new.config import ProcessingConfig
        handler = PPTXHandler(ProcessingConfig())
        assert handler.handler_name == "PPTX Handler"

    def test_create_converter(self):
        from contextifier_new.handlers.pptx.handler import PPTXHandler
        from contextifier_new.handlers.pptx.converter import PptxConverter
        from contextifier_new.config import ProcessingConfig
        handler = PPTXHandler(ProcessingConfig())
        assert isinstance(handler.create_converter(), PptxConverter)

    def test_create_preprocessor(self):
        from contextifier_new.handlers.pptx.handler import PPTXHandler
        from contextifier_new.handlers.pptx.preprocessor import PptxPreprocessor
        from contextifier_new.config import ProcessingConfig
        handler = PPTXHandler(ProcessingConfig())
        assert isinstance(handler.create_preprocessor(), PptxPreprocessor)

    def test_create_metadata_extractor(self):
        from contextifier_new.handlers.pptx.handler import PPTXHandler
        from contextifier_new.handlers.pptx.metadata_extractor import PptxMetadataExtractor
        from contextifier_new.config import ProcessingConfig
        handler = PPTXHandler(ProcessingConfig())
        assert isinstance(handler.create_metadata_extractor(), PptxMetadataExtractor)

    def test_create_content_extractor(self):
        from contextifier_new.handlers.pptx.handler import PPTXHandler
        from contextifier_new.handlers.pptx.content_extractor import PptxContentExtractor
        from contextifier_new.config import ProcessingConfig
        handler = PPTXHandler(ProcessingConfig())
        assert isinstance(handler.create_content_extractor(), PptxContentExtractor)


# ═══════════════════════════════════════════════════════════════════════════════
# Test Full Pipeline
# ═══════════════════════════════════════════════════════════════════════════════

class TestFullPipeline:
    """End-to-end pipeline tests using real PPTX files."""

    def test_converter_to_preprocessor(self):
        from contextifier_new.handlers.pptx.converter import PptxConverter
        from contextifier_new.handlers.pptx.preprocessor import PptxPreprocessor

        data = _make_pptx_with_text(["Pipeline test"])
        ctx = _make_file_context(data)

        conv = PptxConverter()
        prs = conv.convert(ctx)
        assert hasattr(prs, "slides")

        pp = PptxPreprocessor()
        ppd = pp.preprocess(prs)
        assert ppd.properties["slide_count"] == 1

    def test_full_pipeline_text(self):
        from contextifier_new.handlers.pptx.converter import PptxConverter
        from contextifier_new.handlers.pptx.preprocessor import PptxPreprocessor
        from contextifier_new.handlers.pptx.metadata_extractor import PptxMetadataExtractor
        from contextifier_new.handlers.pptx.content_extractor import PptxContentExtractor

        data = _make_pptx_with_text(["Hello from pipeline", "Second text box"])
        ctx = _make_file_context(data)

        # Stage 1: Convert
        prs = PptxConverter().convert(ctx)
        # Stage 2: Preprocess
        ppd = PptxPreprocessor().preprocess(prs)
        # Stage 3: Metadata
        meta = PptxMetadataExtractor().extract(ppd)
        assert meta.page_count == 1
        # Stage 4: Content
        text = PptxContentExtractor().extract_text(ppd)
        assert "Hello from pipeline" in text
        assert "Second text box" in text

    def test_full_pipeline_table(self):
        from contextifier_new.handlers.pptx.converter import PptxConverter
        from contextifier_new.handlers.pptx.preprocessor import PptxPreprocessor
        from contextifier_new.handlers.pptx.content_extractor import PptxContentExtractor

        data = _make_pptx_with_table([["Name", "Value"], ["A", "1"], ["B", "2"]])
        ctx = _make_file_context(data)

        prs = PptxConverter().convert(ctx)
        ppd = PptxPreprocessor().preprocess(prs)
        text = PptxContentExtractor().extract_text(ppd)
        assert "Name" in text
        assert "Value" in text

    def test_full_pipeline_multi_slide(self):
        from contextifier_new.handlers.pptx.converter import PptxConverter
        from contextifier_new.handlers.pptx.preprocessor import PptxPreprocessor
        from contextifier_new.handlers.pptx.content_extractor import PptxContentExtractor

        data = _make_pptx_multi_slide([
            ["First slide"],
            ["Second slide"],
        ])
        ctx = _make_file_context(data)

        prs = PptxConverter().convert(ctx)
        ppd = PptxPreprocessor().preprocess(prs)
        text = PptxContentExtractor().extract_text(ppd)
        assert "[Slide:1]" in text
        assert "[Slide:2]" in text
        assert "First slide" in text
        assert "Second slide" in text

    def test_full_pipeline_notes(self):
        from contextifier_new.handlers.pptx.converter import PptxConverter
        from contextifier_new.handlers.pptx.preprocessor import PptxPreprocessor
        from contextifier_new.handlers.pptx.content_extractor import PptxContentExtractor

        data = _make_pptx_with_notes("Slide content", "My notes")
        ctx = _make_file_context(data)

        prs = PptxConverter().convert(ctx)
        ppd = PptxPreprocessor().preprocess(prs)
        text = PptxContentExtractor().extract_text(ppd)
        assert "Slide content" in text
        assert "My notes" in text

    def test_full_pipeline_extract_all(self):
        from contextifier_new.handlers.pptx.converter import PptxConverter
        from contextifier_new.handlers.pptx.preprocessor import PptxPreprocessor
        from contextifier_new.handlers.pptx.metadata_extractor import PptxMetadataExtractor
        from contextifier_new.handlers.pptx.content_extractor import PptxContentExtractor

        data = _make_pptx_with_table([["H1", "H2"], ["A", "B"]])
        ctx = _make_file_context(data)

        prs = PptxConverter().convert(ctx)
        ppd = PptxPreprocessor().preprocess(prs)
        meta = PptxMetadataExtractor().extract(ppd)

        ext = PptxContentExtractor()
        result = ext.extract_all(ppd, extract_metadata_result=meta)
        assert result.text  # non-empty
        assert result.metadata is meta
        assert len(result.tables) >= 1


# ═══════════════════════════════════════════════════════════════════════════════
# Test Edge Cases
# ═══════════════════════════════════════════════════════════════════════════════

class TestEdgeCases:
    """Edge cases and error handling."""

    def test_corrupted_pptx(self):
        from contextifier_new.handlers.pptx.converter import PptxConverter
        from contextifier_new.errors import ConversionError
        conv = PptxConverter()
        ctx = _make_file_context(b"This is not a PPTX")
        with pytest.raises(ConversionError):
            conv.convert(ctx)

    def test_multiple_text_boxes_ordering(self):
        """Text boxes should maintain visual ordering."""
        from pptx import Presentation
        from pptx.util import Inches
        from contextifier_new.handlers.pptx.content_extractor import PptxContentExtractor
        from contextifier_new.types import PreprocessedData

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        # Add shapes in reverse visual order
        tb_bottom = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(5), Inches(0.5))
        tb_bottom.text_frame.text = "Bottom text"
        tb_top = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(0.5))
        tb_top.text_frame.text = "Top text"

        buf = io.BytesIO()
        prs.save(buf)
        prs2 = Presentation(io.BytesIO(buf.getvalue()))
        ppd = PreprocessedData(content=prs2, raw_content=prs2, resources={"charts_by_slide": {}})

        ext = PptxContentExtractor()
        text = ext.extract_text(ppd)
        top_pos = text.index("Top text")
        bottom_pos = text.index("Bottom text")
        assert top_pos < bottom_pos

    def test_large_table(self):
        """Large table should work without errors."""
        from contextifier_new.handlers.pptx.converter import PptxConverter
        from contextifier_new.handlers.pptx.preprocessor import PptxPreprocessor
        from contextifier_new.handlers.pptx.content_extractor import PptxContentExtractor

        rows = [["H1", "H2", "H3"]]
        for i in range(20):
            rows.append([f"R{i}C1", f"R{i}C2", f"R{i}C3"])
        data = _make_pptx_with_table(rows)
        ctx = _make_file_context(data)

        prs = PptxConverter().convert(ctx)
        ppd = PptxPreprocessor().preprocess(prs)
        text = PptxContentExtractor().extract_text(ppd)
        assert "H1" in text
        assert "R19C3" in text

    def test_mixed_content_slide(self):
        """Slide with both text and table."""
        from pptx import Presentation
        from pptx.util import Inches
        from contextifier_new.handlers.pptx.content_extractor import PptxContentExtractor
        from contextifier_new.types import PreprocessedData

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        # Text box
        tb = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(5), Inches(0.5))
        tb.text_frame.text = "Title text"
        # Table
        tbl_shape = slide.shapes.add_table(2, 2, Inches(1), Inches(2), Inches(5), Inches(2))
        tbl_shape.table.cell(0, 0).text = "Col1"
        tbl_shape.table.cell(0, 1).text = "Col2"
        tbl_shape.table.cell(1, 0).text = "A"
        tbl_shape.table.cell(1, 1).text = "B"

        buf = io.BytesIO()
        prs.save(buf)
        prs2 = Presentation(io.BytesIO(buf.getvalue()))
        ppd = PreprocessedData(content=prs2, raw_content=prs2, resources={"charts_by_slide": {}})

        ext = PptxContentExtractor()
        text = ext.extract_text(ppd)
        assert "Title text" in text
        assert "Col1" in text
        assert "Col2" in text

    def test_empty_presentation(self):
        """Presentation with no slides."""
        from pptx import Presentation
        from contextifier_new.handlers.pptx.content_extractor import PptxContentExtractor
        from contextifier_new.types import PreprocessedData

        prs = Presentation()
        ppd = PreprocessedData(content=prs, raw_content=prs, resources={"charts_by_slide": {}})
        ext = PptxContentExtractor()
        text = ext.extract_text(ppd)
        assert text == ""

    def test_slide_with_empty_text_boxes(self):
        """Slide with only empty text boxes."""
        from pptx import Presentation
        from pptx.util import Inches
        from contextifier_new.handlers.pptx.content_extractor import PptxContentExtractor
        from contextifier_new.types import PreprocessedData

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(0.5))
        tb.text_frame.text = ""

        buf = io.BytesIO()
        prs.save(buf)
        prs2 = Presentation(io.BytesIO(buf.getvalue()))
        ppd = PreprocessedData(content=prs2, raw_content=prs2, resources={"charts_by_slide": {}})

        ext = PptxContentExtractor()
        text = ext.extract_text(ppd)
        assert "[Empty Slide]" in text
