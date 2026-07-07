# tests/unit/raw/test_chart_model.py
"""ChartModel (C3) — chart read/write over real Office files.

Every test works on genuine packages: charts authored by openpyxl
(xlsx-hosted, no caches) and python-pptx (pptx-hosted, full caches +
embedded workbook), plus hand-built packages for the chartEx and
missing-cache corners. Write paths are verified by REOPENING the saved
bytes with python-pptx — proving the rewrite is Office-readable, not
just self-consistent.
"""

from __future__ import annotations

import io
import re
import zipfile

import pytest
from openpyxl import Workbook, load_workbook
from openpyxl.chart import BarChart, LineChart, Reference
from pptx import Presentation
from pptx.chart.data import CategoryChartData, XyChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches

from contextifier.raw.chart import ChartSeriesData, load_chart
from contextifier.raw.opc import OpcPackage, RawUnsupportedError

_CHART_PART_RE = re.compile(r"(?:xl|ppt|word)/charts/chart\d+\.xml$")


def _chart_part_names(pkg: OpcPackage) -> list[str]:
    return sorted(n for n in pkg.part_names if _CHART_PART_RE.fullmatch(n))


# -- builders: real files ------------------------------------------------------


def _xlsx_with_charts() -> bytes:
    """openpyxl workbook hosting a BarChart + a LineChart."""
    wb = Workbook()
    ws = wb.active
    for row in [["Cat", "S1", "S2"], ["A", 1, 4], ["B", 2, 5], ["C", 3, 6]]:
        ws.append(row)
    data = Reference(ws, min_col=2, min_row=1, max_col=3, max_row=4)
    cats = Reference(ws, min_col=1, min_row=2, max_row=4)

    bar = BarChart()
    bar.type = "col"
    bar.title = "Bar Title"
    bar.add_data(data, titles_from_data=True)
    bar.set_categories(cats)
    ws.add_chart(bar, "E5")

    line = LineChart()
    line.title = "Line Title"
    line.add_data(data, titles_from_data=True)
    line.set_categories(cats)
    ws.add_chart(line, "E20")

    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def _pptx_with_charts(specs) -> bytes:
    """One slide + chart per (chart_type, categories, [(name, values)...])."""
    prs = Presentation()
    for chart_type, cats, series in specs:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        chart_data = CategoryChartData()
        chart_data.categories = cats
        for name, values in series:
            chart_data.add_series(name, values)
        slide.shapes.add_chart(
            chart_type, Inches(1), Inches(1), Inches(6), Inches(4), chart_data
        )
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _pptx_charts(data: bytes) -> list:
    """Reopen saved bytes with python-pptx and collect the chart objects."""
    prs = Presentation(io.BytesIO(data))
    return [
        shape.chart for slide in prs.slides for shape in slide.shapes if shape.has_chart
    ]


def _mini_package(parts: dict[str, bytes | str]) -> bytes:
    """Hand-built OPC zip: content types + root rels + given parts."""
    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as z:
        z.writestr(
            "[Content_Types].xml",
            '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
            '<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">'
            '<Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>'
            '<Default Extension="xml" ContentType="application/xml"/>'
            "</Types>",
        )
        z.writestr(
            "_rels/.rels",
            '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
            '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships"/>',
        )
        for name, content in parts.items():
            z.writestr(name, content)
    return buf.getvalue()


_CHARTEX_XML = (
    '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
    '<cx:chartSpace xmlns:cx="http://schemas.microsoft.com/office/drawing/2014/chartex"'
    ' xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
    "<cx:chartData>"
    '<cx:data id="0">'
    '<cx:strDim type="cat"><cx:f>Sheet1!$A$2:$A$4</cx:f>'
    '<cx:lvl ptCount="3"><cx:pt idx="0">A</cx:pt><cx:pt idx="1">B</cx:pt>'
    '<cx:pt idx="2">C</cx:pt></cx:lvl></cx:strDim>'
    '<cx:numDim type="val"><cx:f>Sheet1!$B$2:$B$4</cx:f>'
    '<cx:lvl ptCount="3"><cx:pt idx="0">1</cx:pt><cx:pt idx="2">3</cx:pt>'
    "</cx:lvl></cx:numDim>"
    "</cx:data>"
    "</cx:chartData>"
    "<cx:chart>"
    "<cx:title><cx:tx><cx:rich><a:bodyPr/>"
    "<a:p><a:r><a:t>Funnel Title</a:t></a:r><a:r><a:t></a:t></a:r></a:p>"
    "</cx:rich></cx:tx></cx:title>"
    "<cx:plotArea><cx:plotAreaRegion>"
    '<cx:series layoutId="funnel" uniqueId="{F00D}">'
    "<cx:tx><cx:txData><cx:f>Sheet1!$B$1</cx:f><cx:v>S1</cx:v></cx:txData></cx:tx>"
    '<cx:dataId val="0"/>'
    "</cx:series>"
    "</cx:plotAreaRegion></cx:plotArea>"
    "</cx:chart>"
    "</cx:chartSpace>"
)

_NO_CACHE_CHART_XML = (
    '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
    '<c:chartSpace xmlns:c="http://schemas.openxmlformats.org/drawingml/2006/chart"'
    ' xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
    "<c:chart><c:plotArea><c:barChart>"
    '<c:barDir val="col"/><c:grouping val="clustered"/>'
    # ser 0: openpyxl-style — bare c:f references, no caches anywhere
    "<c:ser>"
    '<c:idx val="0"/><c:order val="0"/>'
    "<c:tx><c:strRef><c:f>Sheet1!$B$1</c:f></c:strRef></c:tx>"
    "<c:cat><c:strRef><c:f>Sheet1!$A$2:$A$4</c:f></c:strRef></c:cat>"
    "<c:val><c:numRef><c:f>Sheet1!$B$2:$B$4</c:f></c:numRef></c:val>"
    "</c:ser>"
    # ser 1: numCache shell with a ptCount but zero c:pt children
    "<c:ser>"
    '<c:idx val="1"/><c:order val="1"/>'
    "<c:val><c:numRef><c:f>Sheet1!$C$2:$C$4</c:f>"
    '<c:numCache><c:ptCount val="3"/></c:numCache>'
    "</c:numRef></c:val>"
    "</c:ser>"
    "</c:barChart></c:plotArea></c:chart>"
    "</c:chartSpace>"
)


# -- 1. openpyxl-hosted reading -------------------------------------------------


class TestOpenpyxlHostedReading:
    def test_kind_title_and_series_shape(self):
        pkg = OpcPackage.open(_xlsx_with_charts())
        names = _chart_part_names(pkg)
        assert len(names) == 2

        models = {m.kind: m for m in (load_chart(pkg, n) for n in names)}
        assert set(models) == {"column", "line"}
        assert models["column"].title == "Bar Title"
        assert models["line"].title == "Line Title"

        for model in models.values():
            series = model.series
            assert len(series) == 2
            # openpyxl writes bare c:f references into the host workbook,
            # never caches — the documented cache-based read is empty.
            for ser in series:
                assert ser.name is None
                assert ser.categories == []
                assert ser.values == []


# -- 2. python-pptx-hosted read + set_data round-trip ---------------------------


class TestPptxSetDataRoundTrip:
    def _package(self) -> tuple[OpcPackage, dict[str, str]]:
        src = _pptx_with_charts(
            [
                (
                    XL_CHART_TYPE.COLUMN_CLUSTERED,
                    ["East", "West", "Mid"],
                    [("Q1", (19.2, 21.4, 16.7)), ("Q2", (22.3, 28.6, 15.2))],
                ),
                (XL_CHART_TYPE.PIE, ["A", "B", "C"], [("Share", (1, 2, 3))]),
            ]
        )
        pkg = OpcPackage.open(src)
        by_kind = {load_chart(pkg, n).kind: n for n in _chart_part_names(pkg)}
        assert set(by_kind) == {"column", "pie"}
        return pkg, by_kind

    def test_reads_pptx_caches(self):
        pkg, by_kind = self._package()
        series = load_chart(pkg, by_kind["column"]).series
        assert [s.name for s in series] == ["Q1", "Q2"]
        assert series[0].categories == ["East", "West", "Mid"]
        assert series[0].values == pytest.approx([19.2, 21.4, 16.7])
        assert series[1].values == pytest.approx([22.3, 28.6, 15.2])

    def test_column_set_data_reopens_in_office_model(self):
        pkg, by_kind = self._package()
        model = load_chart(pkg, by_kind["column"])
        new_cats = ["N", "S", "E", "W"]
        model.set_data(new_cats, [("Alpha", [1, 2, 3, 4]), ("Beta", [5, 6, 7, 8])])
        out = pkg.to_bytes()

        chart = next(
            c
            for c in _pptx_charts(out)
            if c.chart_type == XL_CHART_TYPE.COLUMN_CLUSTERED
        )
        assert list(chart.plots[0].categories) == new_cats
        assert [(s.name, list(s.values)) for s in chart.series] == [
            ("Alpha", [1, 2, 3, 4]),
            ("Beta", [5, 6, 7, 8]),
        ]

    def test_pie_set_data_reopens_in_office_model(self):
        pkg, by_kind = self._package()
        model = load_chart(pkg, by_kind["pie"])
        model.set_data(["Red", "Blue"], [("Split", [60, 40])])
        out = pkg.to_bytes()

        chart = next(c for c in _pptx_charts(out) if c.chart_type == XL_CHART_TYPE.PIE)
        assert list(chart.plots[0].categories) == ["Red", "Blue"]
        assert [(s.name, list(s.values)) for s in chart.series] == [("Split", [60, 40])]

    def test_embedded_workbook_is_regenerated(self):
        pkg, by_kind = self._package()
        model = load_chart(pkg, by_kind["column"])
        workbook_part = model.embedded_workbook_part()
        assert workbook_part is not None
        assert workbook_part.name.endswith(".xlsx")

        model.set_data(["N", "S"], [("Alpha", [1, 2]), ("Beta", [5, 6])])
        out = pkg.to_bytes()

        reopened = OpcPackage.open(out)
        wb_bytes = load_chart(reopened, by_kind["column"]).embedded_workbook_part()
        wb = load_workbook(io.BytesIO(wb_bytes.read()))
        rows = list(wb["Sheet1"].values)
        assert rows == [
            (None, "Alpha", "Beta"),
            ("N", 1, 5),
            ("S", 2, 6),
        ]


# -- 3. set_title round-trip ------------------------------------------------------


class TestSetTitle:
    def test_set_title_roundtrip_via_pptx(self):
        src = _pptx_with_charts(
            [(XL_CHART_TYPE.COLUMN_CLUSTERED, ["A", "B"], [("S", (1, 2))])]
        )
        pkg = OpcPackage.open(src)
        name = _chart_part_names(pkg)[0]
        model = load_chart(pkg, name)
        assert model.title is None  # python-pptx writes no explicit title

        model.set_title("First Draft")  # create path
        model.set_title("Quarterly Revenue")  # replace path
        assert model.title == "Quarterly Revenue"
        out = pkg.to_bytes()

        chart = _pptx_charts(out)[0]
        assert chart.has_title
        assert chart.chart_title.text_frame.text == "Quarterly Revenue"
        assert load_chart(OpcPackage.open(out), name).title == "Quarterly Revenue"


# -- 4. series count growth / shrink ---------------------------------------------


class TestSeriesCountChanges:
    def test_growth_one_to_three(self):
        src = _pptx_with_charts(
            [(XL_CHART_TYPE.COLUMN_CLUSTERED, ["A", "B", "C"], [("Only", (1, 2, 3))])]
        )
        pkg = OpcPackage.open(src)
        name = _chart_part_names(pkg)[0]
        load_chart(pkg, name).set_data(
            ["X", "Y"],
            [("S1", [1, 2]), ("S2", [3, 4]), ("S3", [5, 6])],
        )
        out = pkg.to_bytes()

        chart = _pptx_charts(out)[0]
        assert list(chart.plots[0].categories) == ["X", "Y"]
        assert [(s.name, list(s.values)) for s in chart.series] == [
            ("S1", [1, 2]),
            ("S2", [3, 4]),
            ("S3", [5, 6]),
        ]

    def test_shrink_three_to_one(self):
        src = _pptx_with_charts(
            [
                (
                    XL_CHART_TYPE.COLUMN_CLUSTERED,
                    ["A", "B"],
                    [("S1", (1, 2)), ("S2", (3, 4)), ("S3", (5, 6))],
                )
            ]
        )
        pkg = OpcPackage.open(src)
        name = _chart_part_names(pkg)[0]
        # ChartSeriesData input form (its .categories are ignored).
        load_chart(pkg, name).set_data(
            ["X", "Y"], [ChartSeriesData("Solo", values=[9, 8])]
        )
        out = pkg.to_bytes()

        chart = _pptx_charts(out)[0]
        assert list(chart.plots[0].categories) == ["X", "Y"]
        assert [(s.name, list(s.values)) for s in chart.series] == [("Solo", [9, 8])]


# -- 5. validation + chartEx ------------------------------------------------------


class TestValidationAndChartEx:
    def test_ragged_series_raises_value_error(self):
        src = _pptx_with_charts(
            [(XL_CHART_TYPE.COLUMN_CLUSTERED, ["A", "B"], [("S", (1, 2))])]
        )
        pkg = OpcPackage.open(src)
        model = load_chart(pkg, _chart_part_names(pkg)[0])
        with pytest.raises(ValueError, match="categories"):
            model.set_data(["A", "B", "C"], [("S", [1, 2])])

    def test_chartex_reads_and_write_guards(self):
        data = _mini_package({"ppt/charts/chartEx1.xml": _CHARTEX_XML})
        model = load_chart(OpcPackage.open(data), "ppt/charts/chartEx1.xml")

        assert model.kind == "chartex:funnel"
        assert model.title == "Funnel Title"  # empty a:t run skipped
        series = model.series
        assert len(series) == 1
        assert series[0].name == "S1"
        assert series[0].categories == ["A", "B", "C"]
        assert series[0].values == [1.0, None, 3.0]  # idx 1 is a gap

        with pytest.raises(RawUnsupportedError):
            model.set_data(["A"], [("S", [1])])
        with pytest.raises(RawUnsupportedError):
            model.set_title("New")


# -- 6. byte preservation ----------------------------------------------------------


class TestBytePreservation:
    def test_set_data_touches_only_chart_xml_and_its_workbook(self):
        src = _pptx_with_charts(
            [
                (
                    XL_CHART_TYPE.COLUMN_CLUSTERED,
                    ["A", "B"],
                    [("S1", (1, 2)), ("S2", (3, 4))],
                ),
                (XL_CHART_TYPE.PIE, ["A", "B", "C"], [("Share", (1, 2, 3))]),
            ]
        )
        pkg = OpcPackage.open(src)
        by_kind = {load_chart(pkg, n).kind: n for n in _chart_part_names(pkg)}
        model = load_chart(pkg, by_kind["column"])
        workbook_name = model.embedded_workbook_part().name

        model.set_data(["X", "Y"], [("A", [1, 2])])
        out = pkg.to_bytes()

        before = zipfile.ZipFile(io.BytesIO(src))
        after = zipfile.ZipFile(io.BytesIO(out))
        assert sorted(before.namelist()) == sorted(after.namelist())
        changed = {n for n in before.namelist() if before.read(n) != after.read(n)}
        assert changed == {by_kind["column"], workbook_name}


# -- 7. missing caches --------------------------------------------------------------


class TestMissingCaches:
    def test_no_cache_and_empty_cache_read_safely(self):
        data = _mini_package({"xl/charts/chart1.xml": _NO_CACHE_CHART_XML})
        model = load_chart(OpcPackage.open(data), "xl/charts/chart1.xml")

        assert model.kind == "column"
        series = model.series
        assert len(series) == 2
        # ser 0: refs only, no caches anywhere → everything empty
        assert series[0].name is None
        assert series[0].categories == []
        assert series[0].values == []
        # ser 1: numCache with ptCount=3 and zero pts → three Nones
        assert series[1].values == [None, None, None]


# -- extra: scatter xVal/yVal write path -------------------------------------------


class TestScatterXY:
    def _scatter_pptx(self) -> bytes:
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        chart_data = XyChartData()
        ser = chart_data.add_series("S1")
        ser.add_data_point(1, 2)
        ser.add_data_point(3, 4)
        slide.shapes.add_chart(
            XL_CHART_TYPE.XY_SCATTER,
            Inches(1),
            Inches(1),
            Inches(6),
            Inches(4),
            chart_data,
        )
        buf = io.BytesIO()
        prs.save(buf)
        return buf.getvalue()

    def test_scatter_set_data_writes_xval_yval(self):
        pkg = OpcPackage.open(self._scatter_pptx())
        name = _chart_part_names(pkg)[0]
        model = load_chart(pkg, name)
        assert model.kind == "scatter"
        assert model.series[0].values == [2.0, 4.0]  # yVal
        assert model.series[0].categories == ["1", "3"]  # xVal, as str

        model.set_data([10, 20, 30], [("S1", [1, 2, 3])])
        out = pkg.to_bytes()

        reopened = load_chart(OpcPackage.open(out), name)
        assert reopened.series[0].categories == ["10", "20", "30"]
        assert reopened.series[0].values == [1.0, 2.0, 3.0]
        # all-numeric x categories → numRef inside c:xVal
        ser_el = reopened.xml.find("c:chart/c:plotArea/c:scatterChart/c:ser")
        assert ser_el is not None
        from contextifier.raw.xmlpart import NS

        assert ser_el.find("c:xVal/c:numRef/c:numCache", NS) is not None
        assert ser_el.find("c:yVal/c:numRef/c:numCache", NS) is not None
