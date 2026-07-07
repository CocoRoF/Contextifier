# tests/unit/raw/test_open_raw_integration.py
"""Cross-format integration: the README's raw-layer promises, end to end.

These run the exact user-facing flows — open_raw dispatch, chart editing
through the format models' ``charts`` property, combined cell+chart
edits, and DocumentProcessor.open_raw — with Office-library reopens as
the acceptance check.
"""

from __future__ import annotations

import io

import pytest

from contextifier import DocumentProcessor, open_raw
from contextifier.raw import RawUnsupportedError
from contextifier.raw.docx import DocxRawDocument
from contextifier.raw.pptx import PptxRawDocument
from contextifier.raw.xlsx import XlsxRawDocument


def _xlsx_with_chart() -> bytes:
    from openpyxl import Workbook
    from openpyxl.chart import BarChart, Reference

    wb = Workbook()
    ws = wb.active
    ws.title = "Sales"
    ws.append(["Quarter", "Amount"])
    for row in [["Q1", 120], ["Q2", 135], ["Q3", 150]]:
        ws.append(row)
    chart = BarChart()
    chart.title = "Sales by Quarter"
    chart.add_data(
        Reference(ws, min_col=2, min_row=1, max_row=4), titles_from_data=True
    )
    chart.set_categories(Reference(ws, min_col=1, min_row=2, max_row=4))
    ws.add_chart(chart, "D2")
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def _pptx_with_chart() -> bytes:
    from pptx import Presentation
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE
    from pptx.util import Inches

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    data = CategoryChartData()
    data.categories = ["East", "West"]
    data.add_series("Units", (10, 20))
    slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1), Inches(1), Inches(6), Inches(4), data
    )
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


class TestDispatch:
    def test_open_raw_from_path_and_bytes(self, tmp_path):
        data = _xlsx_with_chart()
        path = tmp_path / "s.xlsx"
        path.write_bytes(data)
        assert isinstance(open_raw(path), XlsxRawDocument)
        assert isinstance(open_raw(data), XlsxRawDocument)  # sniffed
        assert isinstance(open_raw(_pptx_with_chart()), PptxRawDocument)

    def test_docx_dispatch(self):
        from docx import Document

        buf = io.BytesIO()
        Document().save(buf)
        assert isinstance(open_raw(buf.getvalue()), DocxRawDocument)

    def test_unsupported_raises(self):
        import zipfile

        buf = io.BytesIO()
        with zipfile.ZipFile(buf, "w") as z:
            z.writestr("hello.txt", "not an office file")
        with pytest.raises(RawUnsupportedError, match="supported"):
            open_raw(buf.getvalue(), extension="zip")

    def test_document_processor_open_raw(self, tmp_path):
        path = tmp_path / "s.xlsx"
        path.write_bytes(_xlsx_with_chart())
        raw = DocumentProcessor().open_raw(str(path))
        assert isinstance(raw, XlsxRawDocument)


class TestReadmeFlowXlsx:
    def test_cell_plus_chart_edit_roundtrip(self, tmp_path):
        """The README quick-start, verbatim behavior."""
        from openpyxl import load_workbook

        raw = open_raw(_xlsx_with_chart())
        raw.sheets["Sales"].set_cell("B3", 142)
        charts = raw.charts
        assert len(charts) == 1 and charts[0].kind in ("bar", "column")
        assert charts[0].title == "Sales by Quarter"
        charts[0].set_data(
            categories=["Q1", "Q2", "Q3"],
            series=[("Sales", [120, 142, 150])],
        )
        out = tmp_path / "edited.xlsx"
        raw.save(out)

        wb = load_workbook(out)
        assert wb["Sales"]["B3"].value == 142
        # chart part still present and parseable by our own model
        raw2 = open_raw(out)
        assert [round(v) for v in raw2.charts[0].series[0].values] == [120, 142, 150]

    def test_ai_view_still_works_on_edited_file(self, tmp_path):
        """The two views coexist: raw-edit, then extract_text reads it."""
        raw = open_raw(_xlsx_with_chart())
        raw.sheets["Sales"].set_cell("A5", "Q4-new")
        out = tmp_path / "for_extract.xlsx"
        raw.save(out)
        text = DocumentProcessor().extract_text(str(out))
        assert "Q4-new" in text


class TestReadmeFlowPptx:
    def test_chart_edit_via_slide(self, tmp_path):
        from pptx import Presentation

        raw = open_raw(_pptx_with_chart())
        slide = raw.slides[0]
        assert slide.chart_part_names, "chart not discovered on slide"
        chart = slide.charts[0]
        chart.set_data(categories=["North", "South"], series=[("Units", [7, 9])])
        out = tmp_path / "deck.pptx"
        raw.save(out)

        prs = Presentation(str(out))
        found = [s for s in prs.slides[0].shapes if getattr(s, "has_chart", False)]
        assert found
        plot = found[0].chart.plots[0]
        assert list(plot.categories) == ["North", "South"]
        assert [round(v) for v in plot.series[0].values] == [7, 9]
