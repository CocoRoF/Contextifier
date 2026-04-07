"""Create test files for all supported formats for deep testing."""
import os
import json
import csv
import struct
import zipfile
from pathlib import Path

BASE = Path(__file__).parent / "test_files"
BASE.mkdir(exist_ok=True)

# ============================================================
# 1. Plain Text Files
# ============================================================
def create_text_files():
    # Basic TXT
    (BASE / "sample.txt").write_text(
        "This is a sample text file for testing.\n"
        "It has multiple lines.\n"
        "Third line with special chars: é, ñ, ü, 한글, 日本語\n"
        "Fourth line with numbers: 12345\n"
        "Fifth line: end of file.\n",
        encoding="utf-8"
    )

    # Markdown
    (BASE / "sample.md").write_text(
        "# Test Markdown Document\n\n"
        "## Section 1: Introduction\n\n"
        "This is a **bold** and *italic* test document.\n\n"
        "### Subsection 1.1\n\n"
        "- Item 1\n- Item 2\n- Item 3\n\n"
        "## Section 2: Code\n\n"
        "```python\ndef hello():\n    print('Hello World')\n```\n\n"
        "## Section 3: Table\n\n"
        "| Name | Age | City |\n"
        "|------|-----|------|\n"
        "| Alice | 30 | Seoul |\n"
        "| Bob | 25 | Tokyo |\n\n"
        "## Section 4: Conclusion\n\n"
        "End of document.\n",
        encoding="utf-8"
    )

    # Python code file
    (BASE / "sample.py").write_text(
        '"""Sample Python module for testing."""\n\n'
        'import os\nimport sys\n\n\n'
        'class DataProcessor:\n'
        '    """Process data from various sources."""\n\n'
        '    def __init__(self, config: dict):\n'
        '        self.config = config\n'
        '        self._cache = {}\n\n'
        '    def process(self, data: list) -> list:\n'
        '        """Process a list of data items."""\n'
        '        results = []\n'
        '        for item in data:\n'
        '            if item not in self._cache:\n'
        '                self._cache[item] = self._transform(item)\n'
        '            results.append(self._cache[item])\n'
        '        return results\n\n'
        '    def _transform(self, item):\n'
        '        return item.upper() if isinstance(item, str) else item\n',
        encoding="utf-8"
    )

    # JSON
    (BASE / "sample.json").write_text(json.dumps({
        "name": "Contextify Test",
        "version": "1.0.0",
        "settings": {
            "chunk_size": 1000,
            "overlap": 200,
            "formats": ["pdf", "docx", "xlsx"]
        },
        "data": [
            {"id": 1, "value": "first"},
            {"id": 2, "value": "second"},
            {"id": 3, "value": "third"}
        ]
    }, indent=2, ensure_ascii=False), encoding="utf-8")

    # YAML
    (BASE / "sample.yaml").write_text(
        "# Sample YAML Configuration\n"
        "name: Contextify Test\n"
        "version: 1.0.0\n"
        "database:\n"
        "  host: localhost\n"
        "  port: 5432\n"
        "  name: testdb\n"
        "features:\n"
        "  - chunking\n"
        "  - ocr\n"
        "  - metadata\n"
        "settings:\n"
        "  chunk_size: 1000\n"
        "  overlap: 200\n",
        encoding="utf-8"
    )

    # XML
    (BASE / "sample.xml").write_text(
        '<?xml version="1.0" encoding="UTF-8"?>\n'
        '<root>\n'
        '  <item id="1">\n'
        '    <name>First Item</name>\n'
        '    <description>Description of first item</description>\n'
        '  </item>\n'
        '  <item id="2">\n'
        '    <name>Second Item</name>\n'
        '    <description>Description of second item</description>\n'
        '  </item>\n'
        '</root>\n',
        encoding="utf-8"
    )

    # INI
    (BASE / "sample.ini").write_text(
        "[general]\n"
        "name = Contextify\n"
        "version = 0.2.6\n\n"
        "[database]\n"
        "host = localhost\n"
        "port = 5432\n\n"
        "[logging]\n"
        "level = INFO\n"
        "file = app.log\n",
        encoding="utf-8"
    )

    # LOG file
    (BASE / "sample.log").write_text(
        "2024-01-01 10:00:00 INFO Starting application\n"
        "2024-01-01 10:00:01 INFO Loading configuration\n"
        "2024-01-01 10:00:02 WARNING Config file not found, using defaults\n"
        "2024-01-01 10:00:05 INFO Server started on port 8080\n"
        "2024-01-01 10:01:00 ERROR Connection timeout to database\n"
        "2024-01-01 10:01:05 INFO Retry connection successful\n",
        encoding="utf-8"
    )

    # UTF-8 BOM file
    (BASE / "sample_bom.txt").write_bytes(
        b'\xef\xbb\xbf' + "UTF-8 BOM test file\nSecond line\n".encode("utf-8")
    )

    # Korean encoding (EUC-KR)
    (BASE / "sample_euckr.txt").write_bytes(
        "한글 테스트 파일입니다.\n두번째 줄입니다.\n세번째 줄: 가나다라마바사\n".encode("euc-kr")
    )

    # Empty file
    (BASE / "empty.txt").write_text("", encoding="utf-8")

    print("  [OK] Text files created")


# ============================================================
# 2. CSV / TSV Files
# ============================================================
def create_csv_files():
    # Basic CSV
    with open(BASE / "sample.csv", "w", newline="", encoding="utf-8") as f:
        writer = csv.writer(f)
        writer.writerow(["Name", "Age", "City", "Score"])
        writer.writerow(["Alice", "30", "Seoul", "95.5"])
        writer.writerow(["Bob", "25", "Tokyo", "88.3"])
        writer.writerow(["Charlie", "35", "New York", "92.1"])
        writer.writerow(["Diana", "28", "London", "97.8"])
        writer.writerow(["Eve", "22", "Paris", "85.0"])

    # CSV with special characters and commas in fields
    with open(BASE / "complex.csv", "w", newline="", encoding="utf-8") as f:
        writer = csv.writer(f)
        writer.writerow(["Name", "Description", "Price"])
        writer.writerow(["Widget A", "A small, compact widget", "$10.99"])
        writer.writerow(["Widget B", 'Contains "quotes" inside', "$25.50"])
        writer.writerow(["Widget C", "Multi\nline\ndescription", "$15.00"])
        writer.writerow(["한글 제품", "한글 설명, 쉼표 포함", "₩15,000"])

    # Semicolon-delimited CSV
    with open(BASE / "semicolon.csv", "w", newline="", encoding="utf-8") as f:
        f.write("Name;Age;City\n")
        f.write("Alice;30;Seoul\n")
        f.write("Bob;25;Tokyo\n")

    # Large CSV (1000 rows)
    with open(BASE / "large.csv", "w", newline="", encoding="utf-8") as f:
        writer = csv.writer(f)
        writer.writerow(["ID", "Name", "Value", "Category", "Timestamp"])
        for i in range(1000):
            writer.writerow([i, f"Item_{i}", round(i * 1.5, 2), f"Cat_{i % 10}", f"2024-01-{(i%28)+1:02d}"])

    # TSV
    with open(BASE / "sample.tsv", "w", newline="", encoding="utf-8") as f:
        writer = csv.writer(f, delimiter="\t")
        writer.writerow(["Product", "Quantity", "Price", "Total"])
        writer.writerow(["Laptop", "2", "1200.00", "2400.00"])
        writer.writerow(["Mouse", "10", "25.99", "259.90"])
        writer.writerow(["Keyboard", "5", "75.00", "375.00"])

    # Empty CSV
    (BASE / "empty.csv").write_text("", encoding="utf-8")

    # Header-only CSV
    (BASE / "header_only.csv").write_text("Name,Age,City\n", encoding="utf-8")

    print("  [OK] CSV/TSV files created")


# ============================================================
# 3. HTML Files
# ============================================================
def create_html_files():
    # Basic HTML
    (BASE / "sample.html").write_text(
        '<!DOCTYPE html>\n<html lang="en">\n<head>\n'
        '  <meta charset="UTF-8">\n'
        '  <title>Test HTML Document</title>\n'
        '</head>\n<body>\n'
        '  <h1>Main Title</h1>\n'
        '  <p>This is a paragraph with <strong>bold</strong> and <em>italic</em> text.</p>\n'
        '  <h2>Section 1</h2>\n'
        '  <p>Content of section 1.</p>\n'
        '  <ul>\n    <li>Item A</li>\n    <li>Item B</li>\n    <li>Item C</li>\n  </ul>\n'
        '  <h2>Section 2: Table</h2>\n'
        '  <table border="1">\n'
        '    <thead>\n      <tr><th>Name</th><th>Age</th><th>City</th></tr>\n    </thead>\n'
        '    <tbody>\n'
        '      <tr><td>Alice</td><td>30</td><td>Seoul</td></tr>\n'
        '      <tr><td>Bob</td><td>25</td><td>Tokyo</td></tr>\n'
        '    </tbody>\n'
        '  </table>\n'
        '  <h2>Section 3: Code</h2>\n'
        '  <pre><code>function hello() {\n    console.log("Hello!");\n}</code></pre>\n'
        '</body>\n</html>\n',
        encoding="utf-8"
    )

    # HTML with merged table cells (colspan/rowspan)
    (BASE / "complex_table.html").write_text(
        '<!DOCTYPE html>\n<html>\n<body>\n'
        '  <h1>Complex Table Test</h1>\n'
        '  <table border="1">\n'
        '    <tr><th colspan="3">Quarterly Report</th></tr>\n'
        '    <tr><th>Region</th><th>Q1</th><th>Q2</th></tr>\n'
        '    <tr><td rowspan="2">Asia</td><td>100</td><td>120</td></tr>\n'
        '    <tr><td>110</td><td>130</td></tr>\n'
        '    <tr><td>Europe</td><td>200</td><td>210</td></tr>\n'
        '  </table>\n'
        '</body>\n</html>\n',
        encoding="utf-8"
    )

    # HTML with XSS-like content (security test)
    (BASE / "xss_test.html").write_text(
        '<!DOCTYPE html>\n<html>\n<body>\n'
        '  <h1>Security Test</h1>\n'
        '  <p>Normal content here.</p>\n'
        '  <table>\n'
        '    <tr><td>&lt;script&gt;alert("xss")&lt;/script&gt;</td><td>safe</td></tr>\n'
        '    <tr><td>normal</td><td>data</td></tr>\n'
        '  </table>\n'
        '</body>\n</html>\n',
        encoding="utf-8"
    )

    # HTML with Korean content
    (BASE / "korean.html").write_text(
        '<!DOCTYPE html>\n<html lang="ko">\n<head><meta charset="UTF-8"><title>한글 문서</title></head>\n'
        '<body>\n'
        '  <h1>한글 테스트 문서</h1>\n'
        '  <p>이것은 한글로 작성된 테스트 문서입니다.</p>\n'
        '  <table>\n'
        '    <tr><th>이름</th><th>나이</th><th>도시</th></tr>\n'
        '    <tr><td>홍길동</td><td>30</td><td>서울</td></tr>\n'
        '    <tr><td>김철수</td><td>25</td><td>부산</td></tr>\n'
        '  </table>\n'
        '</body>\n</html>\n',
        encoding="utf-8"
    )

    print("  [OK] HTML files created")


# ============================================================
# 4. RTF Files
# ============================================================
def create_rtf_files():
    # Basic RTF
    (BASE / "sample.rtf").write_text(
        r'{\rtf1\ansi\deff0'
        r'{\fonttbl{\f0 Times New Roman;}}'
        r'\pard\f0\fs24 This is a sample RTF document.\par'
        r'It has multiple paragraphs.\par'
        r'\b Bold text\b0  and \i italic text\i0.\par'
        r'Third paragraph with numbers: 12345.\par'
        r'}',
        encoding="utf-8"
    )

    # RTF with table
    (BASE / "table.rtf").write_text(
        r'{\rtf1\ansi\deff0'
        r'{\fonttbl{\f0 Arial;}}'
        r'\pard\f0\fs24 Document with Table\par\par'
        r'\trowd\cellx3000\cellx6000\cellx9000'
        r'\intbl Name\cell Age\cell City\cell\row'
        r'\trowd\cellx3000\cellx6000\cellx9000'
        r'\intbl Alice\cell 30\cell Seoul\cell\row'
        r'\trowd\cellx3000\cellx6000\cellx9000'
        r'\intbl Bob\cell 25\cell Tokyo\cell\row'
        r'\pard After the table.\par'
        r'}',
        encoding="utf-8"
    )

    print("  [OK] RTF files created")


# ============================================================
# 5. DOCX Files (using python-docx)
# ============================================================
def create_docx_files():
    try:
        from docx import Document
        from docx.shared import Inches, Pt
        from docx.enum.text import WD_PARAGRAPH_ALIGNMENT

        # Basic DOCX
        doc = Document()
        doc.add_heading("Test Document", level=0)
        doc.add_paragraph("This is a test DOCX document created for deep testing of Contextify.")
        doc.add_heading("Section 1: Text Content", level=1)
        doc.add_paragraph("Regular paragraph with normal text. This paragraph tests basic text extraction.")
        p = doc.add_paragraph()
        p.add_run("Bold text ").bold = True
        p.add_run("and ")
        r = p.add_run("italic text")
        r.italic = True

        doc.add_heading("Section 2: Lists", level=1)
        doc.add_paragraph("First item", style="List Bullet")
        doc.add_paragraph("Second item", style="List Bullet")
        doc.add_paragraph("Third item", style="List Bullet")

        doc.add_heading("Section 3: Table", level=1)
        table = doc.add_table(rows=4, cols=3)
        table.style = "Table Grid"
        headers = ["Name", "Age", "City"]
        data = [["Alice", "30", "Seoul"], ["Bob", "25", "Tokyo"], ["Charlie", "35", "New York"]]
        for i, h in enumerate(headers):
            table.rows[0].cells[i].text = h
        for r_idx, row_data in enumerate(data):
            for c_idx, val in enumerate(row_data):
                table.rows[r_idx + 1].cells[c_idx].text = val

        doc.add_heading("Section 4: Conclusion", level=1)
        doc.add_paragraph("End of document. Testing metadata extraction as well.")

        doc.core_properties.title = "Deep Test Document"
        doc.core_properties.author = "Contextify Tester"
        doc.core_properties.subject = "Deep Testing"

        doc.save(str(BASE / "sample.docx"))

        # Multi-page DOCX (simulated with many paragraphs)
        doc2 = Document()
        doc2.add_heading("Long Document Test", level=0)
        for i in range(50):
            doc2.add_heading(f"Chapter {i+1}", level=1)
            for j in range(5):
                doc2.add_paragraph(
                    f"This is paragraph {j+1} of chapter {i+1}. "
                    f"Lorem ipsum dolor sit amet, consectetur adipiscing elit. "
                    f"Sed do eiusmod tempor incididunt ut labore et dolore magna aliqua. "
                    f"Value: {i*5+j}"
                )
        doc2.save(str(BASE / "long.docx"))

        # DOCX with merged cells
        doc3 = Document()
        doc3.add_heading("Merged Cell Test", level=0)
        table = doc3.add_table(rows=4, cols=3)
        table.style = "Table Grid"
        table.rows[0].cells[0].text = "Merged Header"
        table.rows[0].cells[0].merge(table.rows[0].cells[2])
        table.rows[1].cells[0].text = "Region"
        table.rows[1].cells[1].text = "Q1"
        table.rows[1].cells[2].text = "Q2"
        table.rows[2].cells[0].text = "Asia"
        table.rows[2].cells[1].text = "100"
        table.rows[2].cells[2].text = "120"
        table.rows[3].cells[0].text = "Europe"
        table.rows[3].cells[1].text = "200"
        table.rows[3].cells[2].text = "210"
        doc3.save(str(BASE / "merged_cells.docx"))

        print("  [OK] DOCX files created")
    except Exception as e:
        print(f"  [FAIL] DOCX: {e}")


# ============================================================
# 6. PPTX Files (using python-pptx)
# ============================================================
def create_pptx_files():
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt

        # Basic PPTX
        prs = Presentation()

        # Slide 1: Title
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "Test Presentation"
        slide.placeholders[1].text = "Created for Contextify Deep Testing"

        # Slide 2: Content
        slide2 = prs.slides.add_slide(prs.slide_layouts[1])
        slide2.shapes.title.text = "Content Slide"
        tf = slide2.placeholders[1].text_frame
        tf.text = "First bullet point"
        p = tf.add_paragraph()
        p.text = "Second bullet point"
        p = tf.add_paragraph()
        p.text = "Third bullet point with 한글"

        # Slide 3: Table
        slide3 = prs.slides.add_slide(prs.slide_layouts[5])  # blank
        slide3.shapes.title.text = "Table Slide"
        table = slide3.shapes.add_table(3, 3, Inches(1), Inches(2), Inches(8), Inches(3)).table
        table.cell(0, 0).text = "Product"
        table.cell(0, 1).text = "Qty"
        table.cell(0, 2).text = "Price"
        table.cell(1, 0).text = "Laptop"
        table.cell(1, 1).text = "5"
        table.cell(1, 2).text = "$1200"
        table.cell(2, 0).text = "Mouse"
        table.cell(2, 1).text = "20"
        table.cell(2, 2).text = "$25"

        # Slide 4: Notes
        slide4 = prs.slides.add_slide(prs.slide_layouts[1])
        slide4.shapes.title.text = "Slide with Notes"
        slide4.placeholders[1].text = "This slide has speaker notes."
        notes_slide = slide4.notes_slide
        notes_slide.notes_text_frame.text = "These are speaker notes for slide 4. Important talking points here."

        prs.save(str(BASE / "sample.pptx"))
        print("  [OK] PPTX files created")
    except Exception as e:
        print(f"  [FAIL] PPTX: {e}")


# ============================================================
# 7. XLSX Files (using openpyxl)
# ============================================================
def create_xlsx_files():
    try:
        from openpyxl import Workbook
        from openpyxl.chart import BarChart, Reference

        # Basic XLSX
        wb = Workbook()
        ws = wb.active
        ws.title = "Sales Data"
        headers = ["Product", "Q1", "Q2", "Q3", "Q4", "Total"]
        ws.append(headers)
        data = [
            ["Laptop", 100, 120, 110, 130, "=SUM(B2:E2)"],
            ["Phone", 200, 180, 220, 250, "=SUM(B3:E3)"],
            ["Tablet", 50, 60, 55, 70, "=SUM(B4:E4)"],
            ["Monitor", 30, 35, 40, 45, "=SUM(B5:E5)"],
        ]
        for row in data:
            ws.append(row)

        # Add a chart
        chart = BarChart()
        chart.type = "col"
        chart.title = "Quarterly Sales"
        chart.x_axis.title = "Product"
        chart.y_axis.title = "Units"
        data_ref = Reference(ws, min_col=2, min_row=1, max_col=5, max_row=5)
        cats = Reference(ws, min_col=1, min_row=2, max_row=5)
        chart.add_data(data_ref, titles_from_data=True)
        chart.set_categories(cats)
        ws.add_chart(chart, "H2")

        # Second sheet
        ws2 = wb.create_sheet("Regions")
        ws2.append(["Region", "Revenue", "Expenses", "Profit"])
        ws2.append(["Asia", 50000, 30000, 20000])
        ws2.append(["Europe", 45000, 28000, 17000])
        ws2.append(["Americas", 60000, 35000, 25000])

        # Merged cells sheet
        ws3 = wb.create_sheet("Merged")
        ws3.merge_cells("A1:C1")
        ws3["A1"] = "Merged Header Row"
        ws3["A2"] = "Col A"
        ws3["B2"] = "Col B"
        ws3["C2"] = "Col C"
        ws3["A3"] = "Data 1"
        ws3["B3"] = "Data 2"
        ws3["C3"] = "Data 3"

        wb.save(str(BASE / "sample.xlsx"))

        # Large XLSX
        wb2 = Workbook()
        ws = wb2.active
        ws.title = "Large Data"
        ws.append(["ID", "Name", "Value", "Category", "Date"])
        for i in range(500):
            ws.append([i, f"Item_{i}", round(i * 2.5, 2), f"Cat_{i%5}", f"2024-{(i%12)+1:02d}-{(i%28)+1:02d}"])
        wb2.save(str(BASE / "large.xlsx"))

        # Empty XLSX
        wb3 = Workbook()
        wb3.save(str(BASE / "empty.xlsx"))

        print("  [OK] XLSX files created")
    except Exception as e:
        print(f"  [FAIL] XLSX: {e}")


# ============================================================
# 8. PDF Files (using reportlab if available, else minimal)
# ============================================================
def create_pdf_files():
    try:
        from reportlab.lib.pagesizes import letter, A4
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, PageBreak
        from reportlab.lib.styles import getSampleStyleSheet
        from reportlab.lib import colors
        from reportlab.lib.units import inch

        styles = getSampleStyleSheet()

        # Basic PDF
        doc = SimpleDocTemplate(str(BASE / "sample.pdf"), pagesize=letter)
        story = []
        story.append(Paragraph("Test PDF Document", styles["Title"]))
        story.append(Spacer(1, 12))
        story.append(Paragraph("This is a test PDF document created for deep testing of the Contextify library.", styles["Normal"]))
        story.append(Spacer(1, 12))
        story.append(Paragraph("Section 1: Introduction", styles["Heading1"]))
        story.append(Paragraph("Contextify processes documents into AI-ready structured text chunks. "
                              "This test validates that PDF extraction works correctly.", styles["Normal"]))
        story.append(Spacer(1, 12))
        story.append(Paragraph("Section 2: Table Data", styles["Heading1"]))

        table_data = [
            ["Name", "Age", "City", "Score"],
            ["Alice", "30", "Seoul", "95.5"],
            ["Bob", "25", "Tokyo", "88.3"],
            ["Charlie", "35", "New York", "92.1"],
        ]
        t = Table(table_data, colWidths=[1.5*inch]*4)
        t.setStyle(TableStyle([
            ("BACKGROUND", (0, 0), (-1, 0), colors.grey),
            ("TEXTCOLOR", (0, 0), (-1, 0), colors.whitesmoke),
            ("GRID", (0, 0), (-1, -1), 1, colors.black),
        ]))
        story.append(t)
        story.append(Spacer(1, 12))
        story.append(Paragraph("Section 3: Conclusion", styles["Heading1"]))
        story.append(Paragraph("End of test document. All sections should be properly extracted.", styles["Normal"]))
        doc.build(story)

        # Multi-page PDF
        doc2 = SimpleDocTemplate(str(BASE / "multipage.pdf"), pagesize=A4)
        story2 = []
        for i in range(10):
            story2.append(Paragraph(f"Page {i+1}: Chapter Title", styles["Title"]))
            story2.append(Spacer(1, 12))
            for j in range(8):
                story2.append(Paragraph(
                    f"Paragraph {j+1} on page {i+1}. Lorem ipsum dolor sit amet, "
                    "consectetur adipiscing elit. Sed do eiusmod tempor incididunt "
                    "ut labore et dolore magna aliqua. Ut enim ad minim veniam.",
                    styles["Normal"]
                ))
                story2.append(Spacer(1, 6))
            story2.append(PageBreak())
        doc2.build(story2)

        # PDF with Korean text
        try:
            from reportlab.pdfbase import pdfmetrics
            from reportlab.pdfbase.ttfonts import TTFont
            # Try registering a CJK font
            doc3 = SimpleDocTemplate(str(BASE / "korean.pdf"), pagesize=A4)
            story3 = []
            story3.append(Paragraph("Korean PDF Test", styles["Title"]))
            story3.append(Paragraph("This PDF contains mixed language content for testing.", styles["Normal"]))
            doc3.build(story3)
        except Exception:
            pass

        print("  [OK] PDF files created")
    except ImportError:
        # Minimal valid PDF without reportlab
        minimal_pdf = b"""%PDF-1.4
1 0 obj
<< /Type /Catalog /Pages 2 0 R >>
endobj
2 0 obj
<< /Type /Pages /Kids [3 0 R] /Count 1 >>
endobj
3 0 obj
<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] /Contents 4 0 R /Resources << /Font << /F1 5 0 R >> >> >>
endobj
4 0 obj
<< /Length 44 >>
stream
BT /F1 12 Tf 100 700 Td (Test PDF) Tj ET
endstream
endobj
5 0 obj
<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>
endobj
xref
0 6
0000000000 65535 f
0000000009 00000 n
0000000058 00000 n
0000000115 00000 n
0000000266 00000 n
0000000360 00000 n
trailer
<< /Size 6 /Root 1 0 R >>
startxref
441
%%EOF"""
        (BASE / "sample.pdf").write_bytes(minimal_pdf)
        print("  [OK] PDF files created (minimal)")
    except Exception as e:
        print(f"  [FAIL] PDF: {e}")


# ============================================================
# 9. XLS Files (minimal binary)
# ============================================================
def create_xls_files():
    try:
        import xlrd
        # xlrd can only read. We'll use xlwt if available
        import xlwt
        wb = xlwt.Workbook()
        ws = wb.add_sheet("Sheet1")
        headers = ["Name", "Age", "City"]
        data = [["Alice", 30, "Seoul"], ["Bob", 25, "Tokyo"]]
        for i, h in enumerate(headers):
            ws.write(0, i, h)
        for r, row in enumerate(data):
            for c, val in enumerate(row):
                ws.write(r+1, c, val)
        wb.save(str(BASE / "sample.xls"))
        print("  [OK] XLS files created")
    except ImportError:
        print("  [SKIP] XLS: xlwt not available")
    except Exception as e:
        print(f"  [FAIL] XLS: {e}")


# ============================================================
# 10. Edge Case Files
# ============================================================
def create_edge_case_files():
    # Very long single line
    (BASE / "long_line.txt").write_text("A" * 100000 + "\n", encoding="utf-8")

    # File with only whitespace
    (BASE / "whitespace.txt").write_text("   \n\n\t\t\n   \n", encoding="utf-8")

    # Binary-looking text
    (BASE / "mixed_binary.txt").write_bytes(
        b"Normal text here\x00\x01\x02More text after nulls\nAnother line\n"
    )

    # Very deeply nested JSON
    nested = {"level": 0}
    current = nested
    for i in range(1, 20):
        current["child"] = {"level": i}
        current = current["child"]
    current["data"] = "deepest value"
    (BASE / "nested.json").write_text(json.dumps(nested, indent=2), encoding="utf-8")

    # CSV with many columns
    with open(BASE / "wide.csv", "w", newline="", encoding="utf-8") as f:
        writer = csv.writer(f)
        writer.writerow([f"Col_{i}" for i in range(100)])
        for r in range(10):
            writer.writerow([f"R{r}_C{i}" for i in range(100)])

    # HTML with deeply nested divs
    nested_html = '<!DOCTYPE html><html><body>'
    for i in range(50):
        nested_html += f'<div class="level-{i}">'
    nested_html += '<p>Deeply nested content</p>'
    for i in range(50):
        nested_html += '</div>'
    nested_html += '</body></html>'
    (BASE / "nested.html").write_text(nested_html, encoding="utf-8")

    print("  [OK] Edge case files created")


# ============================================================
# Main
# ============================================================
if __name__ == "__main__":
    print("Creating test files...")
    create_text_files()
    create_csv_files()
    create_html_files()
    create_rtf_files()
    create_docx_files()
    create_pptx_files()
    create_xlsx_files()
    create_pdf_files()
    create_xls_files()
    create_edge_case_files()

    # List created files
    print(f"\n{'='*60}")
    print(f"Created files in {BASE}:")
    for f in sorted(BASE.iterdir()):
        size = f.stat().st_size
        print(f"  {f.name:30s} {size:>10,d} bytes")
    print(f"{'='*60}")
